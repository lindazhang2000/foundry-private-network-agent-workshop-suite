from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "workshops" / "foundry-private-network-agent-workshop-slide-deck.pptx"


slides_data = [
    {
        "title": "Private Network Azure AI Foundry Agent Workshop",
        "bullets": ["Presenter, date, audience"],
        "notes": "State workshop objective in one sentence.",
    },
    {
        "title": "Why Private Network for Agents",
        "bullets": [
            "Compliance and data boundary requirements",
            "Private endpoint and managed identity patterns",
            "Reduced attack surface",
        ],
        "notes": "Tie to enterprise controls and audit needs.",
        "image": "assets/screenshots/slide-outline/01-keynote-view-6-box.png",
    },
    {
        "title": "End-State Architecture",
        "bullets": [
            "VNet with delegated subnet and PE subnet",
            "AI Services account + model deployment",
            "Project capability host linked to Search, Storage, Cosmos",
        ],
        "notes": "Emphasize two capability host levels: account and project.",
        "image": "assets/screenshots/slide-outline/02-executive-view-slide-friendly.png",
    },
    {
        "title": "Deployment Flow at a Glance",
        "bullets": [
            "Ordered sequence from RG to agent test",
            "Highlight critical dependency points",
        ],
        "notes": "Explain why ordering prevents hard-to-debug failures.",
        "image": "assets/screenshots/slide-outline/03-technical-view-implementation-detail.png",
    },
    {
        "title": "Prerequisites and Guardrails",
        "bullets": [
            "Required roles and permissions",
            "Supported regions and quota checks",
            "CLI and tooling baseline",
        ],
        "notes": "Call out common preflight misses.",
    },
    {
        "title": "Networking Foundation",
        "bullets": [
            "Address spaces and subnet purposes",
            "Delegation requirement for agent-subnet",
        ],
        "notes": "Stress exclusive subnet use for Microsoft.App environments.",
    },
    {
        "title": "Backing Services Security Posture",
        "bullets": [
            "Search, Storage, Cosmos DB hardened settings",
            "Public network access disabled",
        ],
        "notes": "Explain why AAD auth is used for project connections.",
    },
    {
        "title": "Creating AI Services Account via REST",
        "bullets": [
            "Why REST is used (preview API and network injection)",
            "Minimal required payload fields",
        ],
        "notes": "Mention account-level capability host is created from this step.",
    },
    {
        "title": "Private Endpoints and Private DNS",
        "bullets": [
            "Required endpoints and DNS zones",
            "Zone links and A record validation",
        ],
        "notes": "Show how DNS correctness proves private path wiring.",
        "image": "assets/screenshots/slide-outline/04-private-endpoint-evidence.png",
    },
    {
        "title": "Foundry Project and Connections",
        "bullets": [
            "Project identity creation",
            "Connection categories and targets",
            "AAD auth-only configuration",
        ],
        "notes": "Reinforce naming consistency for connection mapping.",
    },
    {
        "title": "RBAC Before Capability Host",
        "bullets": [
            "Storage and Cosmos roles",
            "Search roles",
            "Propagation wait requirement",
        ],
        "notes": "This slide should be treated as non-negotiable gating criteria.",
    },
    {
        "title": "Project Capability Host (Most Critical)",
        "bullets": [
            "Capability host payload fields",
            "Async operation and polling pattern",
            "Success and failure states",
        ],
        "notes": "Explain failure recovery: fix RBAC, recreate the project capability host.",
    },
    {
        "title": "Verification Checklist",
        "bullets": [
            "Capability host status checks",
            "Connection checks",
            "Private endpoint and DNS checks",
            "Model deployment check",
        ],
        "notes": "Ask teams to capture evidence for each check.",
        "image": "assets/screenshots/slide-outline/05-final-validation-checklist.png",
    },
    {
        "title": "Agent Creation and Tool Validation",
        "bullets": [
            "Portal access path (inside private network)",
            "Add Search tool from project connection",
            "Run test prompt",
        ],
        "notes": "Point out this is the final proof of end-to-end success.",
    },
    {
        "title": "Troubleshooting Patterns",
        "bullets": [
            "Invalid endpoint or connection failed",
            "Capability host failed provisioning",
            "Embedding permission denied",
        ],
        "notes": "Map each symptom to likely root cause and fix.",
    },
    {
        "title": "Cleanup and Cost Control",
        "bullets": [
            "Delete project capability host first",
            "Delete resource group",
        ],
        "notes": "Remind teams not to leave resources running.",
    },
    {
        "title": "Recap and Key Takeaways",
        "bullets": [
            "What was built",
            "Why sequence and RBAC matter",
            "Next-step options",
        ],
        "notes": "Provide path to production hardening and automation.",
    },
    {
        "title": "Appendix",
        "bullets": [
            "Command references",
            "API versions used",
            "Role matrix summary",
        ],
        "notes": "Keep appendix in deck for troubleshooting support.",
    },
]


def add_bullets(text_frame, bullets):
    text_frame.clear()
    for idx, bullet in enumerate(bullets):
        p = text_frame.paragraphs[0] if idx == 0 else text_frame.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(20)


def build_deck():
    prs = Presentation()

    for idx, item in enumerate(slides_data):
        if idx == 0:
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = item["title"]
            subtitle = slide.placeholders[1]
            subtitle.text = "Foundry Private Network Agent Workshop Suite"
            title_tf = slide.shapes.title.text_frame
            title_tf.paragraphs[0].font.size = Pt(42)
        else:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = item["title"]
            title_tf = slide.shapes.title.text_frame
            title_tf.paragraphs[0].font.size = Pt(34)

            body = slide.shapes.placeholders[1]
            # Keep room on the right for optional visuals.
            body.left = Inches(0.7)
            body.top = Inches(1.5)
            body.width = Inches(6.1)
            body.height = Inches(5.4)
            add_bullets(body.text_frame, item["bullets"])

            image_rel = item.get("image")
            if image_rel:
                image_path = ROOT / image_rel
                if image_path.exists():
                    slide.shapes.add_picture(str(image_path), Inches(7.0), Inches(1.45), width=Inches(6.05))

        notes = slide.notes_slide.notes_text_frame
        notes.clear()
        notes.text = item["notes"]

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"Created: {OUT}")


if __name__ == "__main__":
    build_deck()
